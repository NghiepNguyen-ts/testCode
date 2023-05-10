
from pptx import Presentation
import numpy as np
from pptx.util import Inches



def finalize_ppt(ppt_path,list_folder,case_name,case_order):
    list_folder.reverse()
    n = 2
    a = list(np.repeat(list_folder, n))
    prs = Presentation(ppt_path)
    slide_count = len(prs.slides)
    left = Inches(10.5)
    top = Inches(1.5)
    width = Inches(1.8)
    height = Inches(1.2)

    for i in range(slide_count):
        slide = prs.slides[i]
        shapes_pic = slide.shapes
        g = "D:\\Yamanami\\Namura\\Code\\Pictures\\" + str(case_order) + "_" + a[i]   + ".png"
        shapes_pic.add_picture(g, left, top, width, height)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[0]
        title_shape.text = 'Report' + ' ' + a[i] 
        prs.save(ppt_path)

## Add slide title
    prs = Presentation(ppt_path)
    title_slide_layout = prs.slide_layouts[7]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Tank Sloshing Report"
    subtitle.text = "Case " + str(case_name)
    prs.save(ppt_path)
## Add slide Overview
    prs = Presentation(ppt_path)
    title_slide_layout = prs.slide_layouts[8]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[0]
    title.text = "Overview"
        # For adjusting the  Margins in inches 
    txBox = slide.placeholders[1]
    # creating textBox
    # creating textFrames
    tf = txBox.text_frame
    tf.text = "Customer: Namura"
    # adding Paragraphs
    p = tf.add_paragraph() 
    # adding text
    p.text = "Purpose: Report the tank sloshing result of case " + str(case_name) 
    p = tf.add_paragraph()
    p.text = "The environment " 
    # ---add table to slide---
    x, y, cx, cy = Inches(1), Inches(3), Inches(6), Inches(1.5)
    shape_table = slide.shapes.add_table(3, 3, x, y, cx, cy)
    table = shape_table.table
    cell_0_0 = table.cell(0,0)
    cell_0_0.text = "Content"
    cell_0_1 = table.cell(0,1)
    cell_0_1.text = "Local machine"
    cell_0_2 = table.cell(0,2)
    cell_0_2.text = "Server"
    cell_1_0 = table.cell(1,0)
    cell_1_0.text = "CPUs/Threads"
    cell_1_1 = table.cell(1,1)
    cell_1_1.text = "8/16"
    cell_1_2 = table.cell(1,2)
    cell_1_2.text = "32/64"
    cell_2_0 = table.cell(2,0)
    cell_2_0.text = "RAM"
    cell_2_1 = table.cell(2,1)
    cell_2_1.text = "48GB"
    cell_2_2 = table.cell(2,2)
    cell_2_2.text = "256GB"
    prs.save(ppt_path)
## Add slide computational setting and running time
    prs = Presentation(ppt_path)
    title_slide_layout = prs.slide_layouts[8]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[0]
    title.text = "Computational setting and running-time"
    txBox = slide.placeholders[1]
    # creating textBox
    # creating textFrames
    tf = txBox.text_frame
    tf.text = "Computational setting and running-time show as the following table, respectively"
    shapes_pic = slide.shapes
    g = "D:\\Yamanami\\Namura\\Code\\Pictures\\" +  str(case_name)   + ".png"
    shapes_pic.add_picture(g, Inches(1), Inches(2.5), Inches(9), Inches(1.5))
    # ---add table to slide---
    x, y, cx, cy = Inches(1), Inches(4.5), Inches(6), Inches(1.5)
    shape_table = slide.shapes.add_table(3, 2, x, y, cx, cy)
    table = shape_table.table
    cell_0_0 = table.cell(0,0)
    cell_0_0.text = "Processing"
    cell_0_1 = table.cell(0,1)
    cell_0_1.text = "Time (s)"
    cell_1_0 = table.cell(1,0)
    cell_1_0.text = "Meshing on local machine"
    cell_2_0 = table.cell(2,0)
    cell_2_0.text = "Solving on the server"
    prs.save(ppt_path)
## Add slide computational setting and running time
    prs = Presentation(ppt_path)
    title_slide_layout = prs.slide_layouts[-1]
    slide = prs.slides.add_slide(title_slide_layout)
    prs.save(ppt_path)